#!/usr/bin/env python3
"""
Génère les 5 templates .pptx squelettes avec design professionnel.
Inspiré des standards McKinsey, BCG, Sequoia, Goldman Sachs.

5 thèmes visuels :
  - corporate    : Bleu marine + blanc, style McKinsey
  - dark         : Fond sombre + accents vifs, style tech
  - finance      : Bleu foncé + doré, style Goldman/JPM
  - startup      : Violet + gradients, style Sequoia/YC
  - minimal      : Noir + blanc + accent unique, style moderne

Usage:
    python build_templates.py [out_dir]
"""
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ─── DESIGN SYSTEM ───────────────────────────────────────────────────────────

THEMES = {
    "corporate": {
        "name": "Corporate McKinsey",
        "primary": (0x00, 0x2B, 0x5C),     # Bleu marine profond
        "accent": (0x00, 0x7E, 0xE5),       # Bleu vif
        "accent2": (0xE6, 0x39, 0x46),       # Rouge accent
        "bg": (0xFF, 0xFF, 0xFF),            # Blanc
        "bg_alt": (0xF7, 0xF8, 0xFA),        # Gris très clair
        "text": (0x1A, 0x1A, 0x1A),          # Quasi-noir
        "text_light": (0x6B, 0x72, 0x80),    # Gris moyen
        "sidebar_width": Inches(0.4),
        "font_title": "Calibri",
        "font_body": "Calibri",
    },
    "dark": {
        "name": "Dark Tech",
        "primary": (0x0D, 0x11, 0x17),       # Fond très sombre
        "accent": (0x58, 0xA6, 0xFF),        # Bleu clair
        "accent2": (0x3F, 0xB9, 0x50),       # Vert
        "bg": (0x0D, 0x11, 0x17),            # Fond sombre
        "bg_alt": (0x16, 0x1B, 0x22),        # Fond sombre alt
        "text": (0xE6, 0xED, 0xF3),          # Blanc cassé
        "text_light": (0x8B, 0x94, 0x9E),    # Gris
        "sidebar_width": Inches(0.35),
        "font_title": "Calibri",
        "font_body": "Calibri",
    },
    "finance": {
        "name": "Finance Premium",
        "primary": (0x1B, 0x2A, 0x4A),       # Bleu nuit
        "accent": (0xC5, 0xA5, 0x5A),        # Doré
        "accent2": (0x2E, 0x86, 0xAB),       # Bleu-vert
        "bg": (0xFF, 0xFF, 0xFF),
        "bg_alt": (0xF5, 0xF3, 0xEF),        # Crème léger
        "text": (0x1A, 0x1A, 0x1A),
        "text_light": (0x6C, 0x75, 0x7D),
        "sidebar_width": Inches(0.35),
        "font_title": "Calibri",
        "font_body": "Calibri",
    },
    "startup": {
        "name": "Startup Sequoia",
        "primary": (0x71, 0x06, 0xEE),       # Violet profond
        "accent": (0xA855, 0xF7, 0x00)[0:3] if False else (0xA8, 0x55, 0xF7),  # Violet clair
        "accent2": (0xEC, 0x48, 0x99),       # Rose
        "bg": (0xFB, 0xF8, 0xFD),            # Lavande très clair
        "bg_alt": (0xF3, 0xEA, 0xFD),        # Lavande
        "text": (0x1A, 0x1A, 0x2E),
        "text_light": (0x59, 0x4F, 0x78),
        "sidebar_width": Inches(0.0),         # Pas de sidebar, style arrondi
        "font_title": "Calibri",
        "font_body": "Calibri",
    },
    "minimal": {
        "name": "Minimal Moderne",
        "primary": (0x11, 0x11, 0x11),       # Noir
        "accent": (0xE6, 0x39, 0x46),        # Rouge accent unique
        "accent2": (0x11, 0x11, 0x11),       # Noir
        "bg": (0xFF, 0xFF, 0xFF),
        "bg_alt": (0xF5, 0xF5, 0xF5),
        "text": (0x11, 0x11, 0x11),
        "text_light": (0x9C, 0xA3, 0xAF),
        "sidebar_width": Inches(0.0),
        "font_title": "Calibri",
        "font_body": "Calibri",
    },
}

TEMPLATES = {
    "executive_deck": {
        "slide_count": 8,
        "theme": "corporate",
        "slides": ["Title", "Executive summary", "Key finding 1", "Key finding 2",
                   "Key finding 3", "Recommandations", "Next steps", "Annexes"],
    },
    "institutional_deck": {
        "slide_count": 22,
        "theme": "corporate",
        "slides": ["Title", "Agenda", "Key message", "Context"] +
                  [f"Section {i}" for i in range(1, 15)] +
                  ["Recos", "Annexes", "Sources"],
    },
    "financial_analysis_deck": {
        "slide_count": 15,
        "theme": "finance",
        "slides": ["Title", "Thèse", "Valorisation", "DCF", "Comps", "Multiples",
                   "Risques", "Catalyseurs", "Scénarios", "Bull case", "Base case",
                   "Bear case", "Recommandation", "Target price", "Sources"],
    },
    "data_deck": {
        "slide_count": 12,
        "theme": "minimal",
        "slides": ["Title", "Key insight"] +
                  [f"Chart {i}" for i in range(1, 9)] +
                  ["Méthodologie", "Dataset"],
    },
    "pitch_deck": {
        "slide_count": 15,
        "theme": "startup",
        "slides": ["Cover", "Problème", "Solution", "Marché (TAM/SAM/SOM)", "Produit",
                   "Traction", "Business model", "GTM", "Compétition", "Équipe",
                   "Roadmap", "Financials", "Ask", "Use of funds", "Contact"],
    },
}

# ─── SLIDE DIMENSIONS ────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _rgb(t):
    """Tuple (r, g, b) -> RGBColor."""
    return RGBColor(t[0], t[1], t[2])


def _add_shape(slide, shape_type, left, top, width, height, fill_rgb, line_rgb=None):
    """Ajouter une forme avec remplissage."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_rgb)
    if line_rgb:
        shape.line.color.rgb = _rgb(line_rgb)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_name, font_size,
                 color, bold=False, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Ajouter une textbox stylée."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    try:
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    except Exception:
        pass
    return tb


def _build_cover_slide(prs, theme, title_text, subtitle_text=""):
    """Slide de couverture avec design impactant."""
    t = THEMES[theme]
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    if theme == "dark":
        # Fond sombre complet
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, t["bg"])
        # Barre accent en haut
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08), t["accent"])
        # Bloc accent à gauche
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(0.5), Inches(2.5), t["accent"])
        txt_color = t["text"]
    elif theme == "finance":
        # Fond bleu nuit couvrant 60% gauche
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(8), SLIDE_H, t["primary"])
        # Barre dorée verticale
        _add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(8), 0, Inches(0.06), SLIDE_H, t["accent"])
        txt_color = (0xFF, 0xFF, 0xFF)
    elif theme == "startup":
        # Fond lavande
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, t["bg"])
        # Grand cercle décoratif violet semi-transparent
        _add_shape(slide, MSO_SHAPE.OVAL, Inches(8.5), Inches(-1), Inches(6), Inches(6), t["primary"])
        # Petit cercle accent
        _add_shape(slide, MSO_SHAPE.OVAL, Inches(11), Inches(5), Inches(3), Inches(3), t["accent"])
        txt_color = t["text"]
    elif theme == "minimal":
        # Ligne rouge horizontale épaisse
        _add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.3), Inches(2), Inches(0.06), t["accent"])
        txt_color = t["text"]
    else:  # corporate
        # Barre latérale bleu marine
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), SLIDE_H, t["primary"])
        # Ligne accent horizontale
        _add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.5), Inches(5), Inches(0.05), t["accent"])
        txt_color = t["text"]

    # Titre principal
    title_left = Inches(1.2) if theme != "minimal" else Inches(1)
    title_top = Inches(2.0) if theme not in ("finance",) else Inches(2.2)
    _add_textbox(slide, title_left, title_top, Inches(9), Inches(1.5),
                 title_text, t["font_title"], Pt(36), txt_color, bold=True)

    # Sous-titre
    if subtitle_text:
        sub_color = t["text_light"] if theme != "dark" else t["text_light"]
        if theme == "finance":
            sub_color = (0xC5, 0xA5, 0x5A)
        _add_textbox(slide, title_left, Inches(3.8), Inches(9), Inches(0.8),
                     subtitle_text, t["font_body"], Pt(18), sub_color)

    return slide


def _build_content_slide(prs, theme, slide_num, total, title_text):
    """Slide de contenu standard avec tous les éléments de design."""
    t = THEMES[theme]
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    if theme == "dark":
        # Fond sombre
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, t["bg"])

    # ── BARRE D'ACCENT LATÉRALE (corporate, finance) ──
    if t["sidebar_width"] > 0:
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   0, 0, t["sidebar_width"], SLIDE_H, t["primary"])

    # ── LIGNE SUPÉRIEURE ACCENT ──
    line_left = t["sidebar_width"] if t["sidebar_width"] > 0 else 0
    accent_color = t["accent"]
    if theme == "finance":
        accent_color = t["accent"]  # doré
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               line_left, 0, Emu(SLIDE_W - line_left), Inches(0.04), accent_color)

    # ── BARRE DE TITRE (fond légèrement teinté) ──
    title_bar_left = t["sidebar_width"] if t["sidebar_width"] > 0 else Inches(0)
    title_bar_w = Emu(SLIDE_W - title_bar_left)
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               title_bar_left, Inches(0.04), title_bar_w, Inches(1.1),
               t["bg_alt"] if theme != "dark" else t["bg_alt"])

    # ── TITRE (Action title) ──
    title_left = Inches(0.9) if t["sidebar_width"] > 0 else Inches(0.7)
    _add_textbox(slide, title_left, Inches(0.18), Inches(11.5), Inches(0.8),
                 title_text, t["font_title"], Pt(22),
                 t["primary"] if theme != "dark" else t["text"],
                 bold=True)

    # ── SÉPARATEUR SOUS TITRE ──
    sep_y = Inches(1.14)
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               title_left, sep_y, Inches(11), Inches(0.02),
               t["accent"] if theme in ("corporate", "startup") else t["text_light"])

    # ── FOOTER ──
    footer_y = Inches(7.05)
    # Ligne footer
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               title_left, footer_y, Inches(11.5), Inches(0.01),
               t["text_light"] if theme != "dark" else (0x30, 0x36, 0x3D))

    # Numéro de slide
    _add_textbox(slide, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
                 f"{slide_num}/{total}", t["font_body"], Pt(9),
                 t["text_light"], alignment=PP_ALIGN.RIGHT)

    # Placeholder "CONFIDENTIEL" à gauche footer
    _add_textbox(slide, title_left, Inches(7.1), Inches(3), Inches(0.3),
                 "CONFIDENTIEL", t["font_body"], Pt(8), t["text_light"])

    return slide


def _build_section_slide(prs, theme, section_title):
    """Slide de transition entre sections."""
    t = THEMES[theme]
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    if theme == "dark":
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, t["bg"])
        # Barre accent épaisse
        _add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.2),
                   Inches(1.5), Inches(0.06), t["accent"])
    elif theme == "finance":
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, t["primary"])
        # Ligne dorée
        _add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.2),
                   Inches(2.5), Inches(0.04), t["accent"])
    elif theme == "startup":
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, t["bg_alt"])
        # Cercle décoratif
        _add_shape(slide, MSO_SHAPE.OVAL, Inches(10), Inches(4.5),
                   Inches(4), Inches(4), t["primary"])
    elif theme == "minimal":
        # Numéro de section géant
        _add_textbox(slide, Inches(1), Inches(1), Inches(4), Inches(4),
                     "§", t["font_title"], Pt(120), t["bg_alt"], bold=True)
    else:  # corporate
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, t["primary"])
        # Ligne accent
        _add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.2),
                   Inches(3), Inches(0.05), t["accent"])

    # Titre de section centré
    txt_color = (0xFF, 0xFF, 0xFF) if theme in ("corporate", "finance") else t["text"]
    if theme == "dark":
        txt_color = t["text"]
    _add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
                 section_title, t["font_title"], Pt(32), txt_color,
                 bold=True, alignment=PP_ALIGN.LEFT)

    return slide


def build_template(name, spec, out_dir):
    """Construit un template .pptx professionnel."""
    if not HAS_PPTX:
        print(f"[SKIP {name}] python-pptx non installé")
        md = out_dir / f"{name}.md"
        md.write_text(f"# Template {name}\n\n" +
                      "\n".join(f"- Slide {i+1}: {s}" for i, s in enumerate(spec["slides"])),
                      encoding="utf-8")
        return

    theme = spec.get("theme", "corporate")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = spec["slide_count"]

    for i, title in enumerate(spec["slides"]):
        if i == 0:
            _build_cover_slide(prs, theme, f"[{name}]", f"Template — {total} slides")
        elif title.startswith("Section") or title == "Agenda" or title == "Context":
            _build_section_slide(prs, theme, title)
        else:
            _build_content_slide(prs, theme, i + 1, total, f"[{i+1}/{total}] {title}")

    out = out_dir / f"{name}.pptx"
    prs.save(str(out))
    print(f"[OK] {out}  (thème: {theme})")


def main():
    out_dir = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent.parent / "templates"
    out_dir.mkdir(parents=True, exist_ok=True)
    for name, spec in TEMPLATES.items():
        build_template(name, spec, out_dir)
    print(f"\n[OK] {len(TEMPLATES)} templates generes dans {out_dir}")


if __name__ == "__main__":
    main()
