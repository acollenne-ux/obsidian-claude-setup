#!/usr/bin/env python3
"""
Moteur de rendu .pptx professionnel — Design McKinsey/BCG/Goldman.

Génère un .pptx éditable à partir d'un YAML de deck avec design premium.
5 thèmes disponibles : corporate, dark, finance, startup, minimal.

Usage:
    python pptx_builder.py deck.yaml out.pptx [--theme corporate|dark|finance|startup|minimal]

Format YAML attendu :
    theme: corporate          # optionnel, défaut=corporate
    title: "Titre principal"
    subtitle: "Sous-titre"
    confidential: true        # optionnel, affiche "CONFIDENTIEL" en footer
    slides:
      - type: cover           # cover | section | content | two_col | kpi | chart | closing
        action_title: "..."
        subtitle: "..."
      - type: content
        action_title: "Le CA progresse de 14% en APAC"
        bullets:
          - "Point clé 1"
          - "Point clé 2"
        source: "Source : Rapport annuel 2025"
        notes: "Notes pour le présentateur"
      - type: two_col
        action_title: "Comparaison avant/après"
        left_title: "Avant"
        left_bullets: ["Point 1", "Point 2"]
        right_title: "Après"
        right_bullets: ["Point 1", "Point 2"]
      - type: kpi
        action_title: "KPIs clés Q3 2025"
        kpis:
          - label: "Revenue"
            value: "$42M"
            delta: "+14%"
          - label: "EBITDA"
            value: "$8.2M"
            delta: "+22%"
      - type: section
        action_title: "Titre de section"
      - type: closing
        action_title: "Merci"
        subtitle: "Questions ?"
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════════════════════════════════
# DESIGN SYSTEM — 5 thèmes professionnels
# ═══════════════════════════════════════════════════════════════════════════════

THEMES = {
    "corporate": {
        "name": "Corporate McKinsey",
        "primary":    (0x00, 0x2B, 0x5C),
        "accent":     (0x00, 0x7E, 0xE5),
        "accent2":    (0xE6, 0x39, 0x46),
        "bg":         (0xFF, 0xFF, 0xFF),
        "bg_alt":     (0xF7, 0xF8, 0xFA),
        "text":       (0x1A, 0x1A, 0x1A),
        "text_light": (0x6B, 0x72, 0x80),
        "sidebar_w":  Inches(0.4),
        "font_title": "Calibri",
        "font_body":  "Calibri",
        "chart_colors": [(0x00,0x2B,0x5C),(0x00,0x7E,0xE5),(0xE6,0x39,0x46),(0x6B,0x72,0x80),(0x34,0xD3,0x99)],
    },
    "dark": {
        "name": "Dark Tech",
        "primary":    (0x0D, 0x11, 0x17),
        "accent":     (0x58, 0xA6, 0xFF),
        "accent2":    (0x3F, 0xB9, 0x50),
        "bg":         (0x0D, 0x11, 0x17),
        "bg_alt":     (0x16, 0x1B, 0x22),
        "text":       (0xE6, 0xED, 0xF3),
        "text_light": (0x8B, 0x94, 0x9E),
        "sidebar_w":  Inches(0.35),
        "font_title": "Calibri",
        "font_body":  "Calibri",
        "chart_colors": [(0x58,0xA6,0xFF),(0x3F,0xB9,0x50),(0xF0,0x88,0x3E),(0xBC,0x3F,0xBC),(0xE6,0xED,0xF3)],
    },
    "finance": {
        "name": "Finance Premium",
        "primary":    (0x1B, 0x2A, 0x4A),
        "accent":     (0xC5, 0xA5, 0x5A),
        "accent2":    (0x2E, 0x86, 0xAB),
        "bg":         (0xFF, 0xFF, 0xFF),
        "bg_alt":     (0xF5, 0xF3, 0xEF),
        "text":       (0x1A, 0x1A, 0x1A),
        "text_light": (0x6C, 0x75, 0x7D),
        "sidebar_w":  Inches(0.35),
        "font_title": "Calibri",
        "font_body":  "Calibri",
        "chart_colors": [(0x1B,0x2A,0x4A),(0xC5,0xA5,0x5A),(0x2E,0x86,0xAB),(0x6C,0x75,0x7D),(0x8B,0x5C,0xF6)],
    },
    "startup": {
        "name": "Startup Sequoia",
        "primary":    (0x71, 0x06, 0xEE),
        "accent":     (0xA8, 0x55, 0xF7),
        "accent2":    (0xEC, 0x48, 0x99),
        "bg":         (0xFB, 0xF8, 0xFD),
        "bg_alt":     (0xF3, 0xEA, 0xFD),
        "text":       (0x1A, 0x1A, 0x2E),
        "text_light": (0x59, 0x4F, 0x78),
        "sidebar_w":  Inches(0),
        "font_title": "Calibri",
        "font_body":  "Calibri",
        "chart_colors": [(0x71,0x06,0xEE),(0xA8,0x55,0xF7),(0xEC,0x48,0x99),(0x06,0xB6,0xD4),(0x22,0xC5,0x5E)],
    },
    "minimal": {
        "name": "Minimal Moderne",
        "primary":    (0x11, 0x11, 0x11),
        "accent":     (0xE6, 0x39, 0x46),
        "accent2":    (0x11, 0x11, 0x11),
        "bg":         (0xFF, 0xFF, 0xFF),
        "bg_alt":     (0xF5, 0xF5, 0xF5),
        "text":       (0x11, 0x11, 0x11),
        "text_light": (0x9C, 0xA3, 0xAF),
        "sidebar_w":  Inches(0),
        "font_title": "Calibri",
        "font_body":  "Calibri",
        "chart_colors": [(0x11,0x11,0x11),(0xE6,0x39,0x46),(0x9C,0xA3,0xAF),(0x4B,0x55,0x63),(0xD1,0xD5,0xDB)],
    },
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _rgb(t):
    return RGBColor(t[0], t[1], t[2])


def _rect(slide, left, top, w, h, fill, line=None):
    """Rectangle rempli, sans bordure par défaut."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(fill)
    if line:
        s.line.color.rgb = _rgb(line)
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def _oval(slide, left, top, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(fill)
    s.line.fill.background()
    return s


def _textbox(slide, left, top, w, h, text, font, size, color,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=None):
    """Textbox avec contrôle fin du style."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    try:
        p.space_before = Pt(0)
        p.space_after = Pt(0)
    except Exception:
        pass
    return tb


def _add_bullets(slide, left, top, w, h, bullets, t, numbered=False):
    """Ajouter des bullets stylées avec espacement généreux."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(6)
        p.space_after = Pt(4)
        p.line_spacing = Pt(24)
        run = p.add_run()
        if numbered:
            prefix = f"{i+1}.  "
        else:
            prefix = "    "  # espace propre, le bullet sera un accent shape
        run.text = f"{prefix}{bullet}"
        run.font.name = t["font_body"]
        run.font.size = Pt(15)
        run.font.color.rgb = _rgb(t["text"])

    # Petits carrés/ronds d'accent pour chaque bullet
    if not numbered:
        for i in range(len(bullets)):
            bullet_y = top + Inches(0.08) + Inches(0.38) * i
            if bullet_y < top + h - Inches(0.2):
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left, bullet_y + Inches(0.05),
                    Inches(0.12), Inches(0.12),
                    )
                dot.fill.solid()
                dot.fill.fore_color.rgb = _rgb(t["accent"])
                dot.line.fill.background()

    return tb


# ═══════════════════════════════════════════════════════════════════════════════
# COMMON ELEMENTS (appliqués à chaque slide de contenu)
# ═══════════════════════════════════════════════════════════════════════════════

def _apply_base(slide, t, theme_name):
    """Appliquer fond + barre latérale + ligne accent top."""
    # Fond (pour dark theme)
    if theme_name == "dark":
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, t["bg"])
    elif theme_name == "startup":
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, t["bg"])

    # Barre latérale
    if t["sidebar_w"] and t["sidebar_w"] > 0:
        _rect(slide, 0, 0, t["sidebar_w"], SLIDE_H, t["primary"])

    # Ligne accent en haut
    line_left = t["sidebar_w"] if t["sidebar_w"] and t["sidebar_w"] > 0 else 0
    _rect(slide, line_left, 0, Emu(SLIDE_W - line_left), Inches(0.045), t["accent"])


def _apply_title_bar(slide, t, theme_name, action_title):
    """Barre de titre avec fond teinté + action title + séparateur."""
    bar_left = t["sidebar_w"] if t["sidebar_w"] and t["sidebar_w"] > 0 else 0
    bar_w = Emu(SLIDE_W - bar_left)

    # Fond de la barre de titre
    bg_fill = t["bg_alt"] if theme_name != "dark" else t["bg_alt"]
    _rect(slide, bar_left, Inches(0.045), bar_w, Inches(1.05), bg_fill)

    # Action title
    title_left = Inches(0.9) if t["sidebar_w"] and t["sidebar_w"] > 0 else Inches(0.7)
    _textbox(slide, title_left, Inches(0.2), Inches(11.5), Inches(0.75),
             action_title, t["font_title"], Pt(21),
             t["primary"] if theme_name != "dark" else t["text"],
             bold=True)

    # Séparateur sous titre
    sep_color = t["accent"] if theme_name in ("corporate", "startup", "finance") else t["text_light"]
    _rect(slide, title_left, Inches(1.08), Inches(11.4), Inches(0.02), sep_color)

    return title_left


def _apply_footer(slide, t, theme_name, slide_num, total, confidential=True):
    """Footer avec ligne + numéro + mention confidentiel."""
    title_left = Inches(0.9) if t["sidebar_w"] and t["sidebar_w"] > 0 else Inches(0.7)
    footer_y = Inches(7.0)

    # Ligne footer
    line_color = t["text_light"] if theme_name != "dark" else (0x30, 0x36, 0x3D)
    _rect(slide, title_left, footer_y, Inches(11.6), Inches(0.012), line_color)

    # Numéro
    _textbox(slide, Inches(12.2), Inches(7.08), Inches(0.8), Inches(0.3),
             f"{slide_num}", t["font_body"], Pt(9), t["text_light"],
             align=PP_ALIGN.RIGHT)

    # Confidentiel
    if confidential:
        _textbox(slide, title_left, Inches(7.08), Inches(3), Inches(0.3),
                 "CONFIDENTIEL", t["font_body"], Pt(8), t["text_light"])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def _slide_cover(prs, t, tn, data):
    """Slide de couverture — design impactant."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = data.get("action_title", data.get("title", "Titre"))
    subtitle = data.get("subtitle", "")

    if tn == "dark":
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, t["bg"])
        _rect(slide, 0, 0, SLIDE_W, Inches(0.07), t["accent"])
        _rect(slide, 0, Inches(2.3), Inches(0.5), Inches(3), t["accent"])
        txt_c = t["text"]
        sub_c = t["text_light"]
    elif tn == "finance":
        _rect(slide, 0, 0, Inches(8.5), SLIDE_H, t["primary"])
        _rect(slide, Inches(8.5), 0, Inches(0.05), SLIDE_H, t["accent"])
        # Motif décoratif doré
        _rect(slide, Inches(0.8), Inches(5.5), Inches(1.5), Inches(0.04), t["accent"])
        txt_c = (0xFF, 0xFF, 0xFF)
        sub_c = t["accent"]
    elif tn == "startup":
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, t["bg"])
        _oval(slide, Inches(9), Inches(-1.5), Inches(6), Inches(6), t["primary"])
        _oval(slide, Inches(10.5), Inches(5), Inches(4), Inches(4), t["accent"])
        _rect(slide, Inches(1), Inches(5.2), Inches(2), Inches(0.05), t["primary"])
        txt_c = t["text"]
        sub_c = t["text_light"]
    elif tn == "minimal":
        _rect(slide, Inches(0.8), Inches(3.4), Inches(2.5), Inches(0.05), t["accent"])
        txt_c = t["text"]
        sub_c = t["text_light"]
    else:  # corporate
        _rect(slide, 0, 0, Inches(0.5), SLIDE_H, t["primary"])
        _rect(slide, Inches(0.5), Inches(3.6), Inches(5), Inches(0.04), t["accent"])
        # Petit accent carré
        _rect(slide, Inches(0.5), Inches(1.8), Inches(0.12), Inches(0.8), t["accent"])
        txt_c = t["text"]
        sub_c = t["text_light"]

    # Titre
    tl = Inches(1.2) if tn not in ("minimal",) else Inches(0.8)
    tt = Inches(2.0) if tn != "finance" else Inches(2.2)
    _textbox(slide, tl, tt, Inches(9), Inches(1.5), title,
             t["font_title"], Pt(36), txt_c, bold=True)

    # Sous-titre
    if subtitle:
        _textbox(slide, tl, Inches(3.8) if tn != "finance" else Inches(4.0),
                 Inches(9), Inches(0.8), subtitle,
                 t["font_body"], Pt(17), sub_c, italic=(tn == "minimal"))

    # Date en bas
    import datetime
    date_str = datetime.date.today().strftime("%B %Y")
    date_y = Inches(6.5)
    _textbox(slide, tl, date_y, Inches(4), Inches(0.3), date_str,
             t["font_body"], Pt(11), t["text_light"])

    return slide


def _slide_section(prs, t, tn, data):
    """Slide de transition entre sections."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = data.get("action_title", "Section")

    if tn == "dark":
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, t["bg"])
        _rect(slide, Inches(1.2), Inches(3.5), Inches(2), Inches(0.05), t["accent"])
        txt_c = t["text"]
    elif tn == "finance":
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, t["primary"])
        _rect(slide, Inches(1.2), Inches(4.2), Inches(2.5), Inches(0.04), t["accent"])
        txt_c = (0xFF, 0xFF, 0xFF)
    elif tn == "startup":
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, t["bg_alt"])
        _oval(slide, Inches(10), Inches(4), Inches(5), Inches(5), t["primary"])
        txt_c = t["text"]
    elif tn == "minimal":
        # Grande lettre décorative
        _textbox(slide, Inches(9), Inches(0.5), Inches(4), Inches(6),
                 "»", t["font_title"], Pt(200), t["bg_alt"], bold=True,
                 align=PP_ALIGN.RIGHT)
        _rect(slide, Inches(1.2), Inches(4.0), Inches(1.5), Inches(0.05), t["accent"])
        txt_c = t["text"]
    else:  # corporate
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, t["primary"])
        _rect(slide, Inches(1.2), Inches(4.2), Inches(3.5), Inches(0.04), t["accent"])
        txt_c = (0xFF, 0xFF, 0xFF)

    _textbox(slide, Inches(1.2), Inches(2.5), Inches(10), Inches(1.5),
             title, t["font_title"], Pt(32), txt_c, bold=True)

    # Sous-titre de section si présent
    sub = data.get("subtitle", "")
    if sub:
        sub_c = t["text_light"] if tn not in ("corporate", "finance") else (0xCC, 0xCC, 0xCC)
        _textbox(slide, Inches(1.2), Inches(4.5), Inches(10), Inches(1),
                 sub, t["font_body"], Pt(16), sub_c)

    return slide


def _slide_content(prs, t, tn, data, slide_num, total, confidential):
    """Slide de contenu standard — action title + bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    action_title = data.get("action_title", "Action title manquant")

    _apply_base(slide, t, tn)
    title_left = _apply_title_bar(slide, t, tn, action_title)
    _apply_footer(slide, t, tn, slide_num, total, confidential)

    # Corps : bullets
    bullets = data.get("bullets", [])
    if bullets:
        body_left = Inches(1.0) if t["sidebar_w"] and t["sidebar_w"] > 0 else Inches(0.8)
        _add_bullets(slide, body_left, Inches(1.4), Inches(11), Inches(5.2),
                     bullets, t, numbered=data.get("numbered", False))

    # Source en bas
    source = data.get("source", "")
    if source:
        _textbox(slide, title_left, Inches(6.6), Inches(10), Inches(0.3),
                 source, t["font_body"], Pt(9), t["text_light"], italic=True)

    # Notes orateur
    notes = data.get("notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def _slide_two_col(prs, t, tn, data, slide_num, total, confidential):
    """Slide en 2 colonnes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    action_title = data.get("action_title", "Comparaison")

    _apply_base(slide, t, tn)
    title_left = _apply_title_bar(slide, t, tn, action_title)
    _apply_footer(slide, t, tn, slide_num, total, confidential)

    col_left = Inches(0.9) if t["sidebar_w"] and t["sidebar_w"] > 0 else Inches(0.7)
    col_w = Inches(5.5)
    body_top = Inches(1.5)

    # Séparateur vertical central
    sep_x = Inches(6.5)
    _rect(slide, sep_x, Inches(1.3), Inches(0.015), Inches(5.2), t["text_light"])

    # Colonne gauche
    left_title = data.get("left_title", "")
    if left_title:
        _textbox(slide, col_left, body_top, col_w, Inches(0.4),
                 left_title, t["font_title"], Pt(16), t["accent"], bold=True)
    left_bullets = data.get("left_bullets", [])
    if left_bullets:
        _add_bullets(slide, col_left, Inches(2.0), col_w, Inches(4.5), left_bullets, t)

    # Colonne droite
    right_left = Inches(7.0)
    right_title = data.get("right_title", "")
    if right_title:
        _textbox(slide, right_left, body_top, col_w, Inches(0.4),
                 right_title, t["font_title"], Pt(16), t["accent2"], bold=True)
    right_bullets = data.get("right_bullets", [])
    if right_bullets:
        _add_bullets(slide, right_left, Inches(2.0), col_w, Inches(4.5), right_bullets, t)

    # Source
    source = data.get("source", "")
    if source:
        _textbox(slide, title_left, Inches(6.6), Inches(10), Inches(0.3),
                 source, t["font_body"], Pt(9), t["text_light"], italic=True)

    return slide


def _slide_kpi(prs, t, tn, data, slide_num, total, confidential):
    """Slide KPI avec cartes métriques."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    action_title = data.get("action_title", "KPIs clés")

    _apply_base(slide, t, tn)
    title_left = _apply_title_bar(slide, t, tn, action_title)
    _apply_footer(slide, t, tn, slide_num, total, confidential)

    kpis = data.get("kpis", [])
    if not kpis:
        return slide

    # Disposition des cartes KPI
    n = len(kpis)
    max_per_row = min(n, 4)
    card_w = Inches(2.8)
    card_h = Inches(2.2)
    gap = Inches(0.4)
    total_w = max_per_row * card_w + (max_per_row - 1) * gap
    start_x = Emu((SLIDE_W - total_w) / 2)
    start_y = Inches(2.2)

    for i, kpi in enumerate(kpis):
        row = i // max_per_row
        col = i % max_per_row
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        # Carte fond
        card_bg = t["bg_alt"] if tn != "dark" else t["bg_alt"]
        card = _rect(slide, x, y, card_w, card_h, card_bg)

        # Barre accent en haut de la carte
        _rect(slide, x, y, card_w, Inches(0.05), t["accent"])

        # Label
        _textbox(slide, x + Inches(0.25), y + Inches(0.25), Inches(2.3), Inches(0.35),
                 kpi.get("label", ""), t["font_body"], Pt(12), t["text_light"],
                 bold=True)

        # Valeur (grande)
        _textbox(slide, x + Inches(0.25), y + Inches(0.7), Inches(2.3), Inches(0.7),
                 kpi.get("value", ""), t["font_title"], Pt(28), t["text"],
                 bold=True)

        # Delta (couleur conditionnelle)
        delta = kpi.get("delta", "")
        if delta:
            delta_color = (0x22, 0xC5, 0x5E) if delta.startswith("+") else (0xEF, 0x44, 0x44)
            if delta.startswith("=") or delta.startswith("~"):
                delta_color = t["text_light"]
            _textbox(slide, x + Inches(0.25), y + Inches(1.5), Inches(2.3), Inches(0.4),
                     delta, t["font_body"], Pt(14), delta_color, bold=True)

    # Source
    source = data.get("source", "")
    if source:
        _textbox(slide, title_left, Inches(6.6), Inches(10), Inches(0.3),
                 source, t["font_body"], Pt(9), t["text_light"], italic=True)

    return slide


def _slide_closing(prs, t, tn, data):
    """Slide de conclusion/remerciement."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = data.get("action_title", "Merci")
    subtitle = data.get("subtitle", "")

    if tn == "dark":
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, t["bg"])
        _rect(slide, 0, Inches(3.4), SLIDE_W, Inches(0.06), t["accent"])
        txt_c = t["text"]
    elif tn == "finance":
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, t["primary"])
        _rect(slide, Inches(5.5), Inches(4.5), Inches(2.5), Inches(0.04), t["accent"])
        txt_c = (0xFF, 0xFF, 0xFF)
    elif tn == "startup":
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, t["bg"])
        _oval(slide, Inches(-2), Inches(-2), Inches(6), Inches(6), t["primary"])
        _oval(slide, Inches(11), Inches(5), Inches(4), Inches(4), t["accent"])
        txt_c = t["text"]
    elif tn == "minimal":
        _rect(slide, Inches(5.5), Inches(4.0), Inches(2.5), Inches(0.05), t["accent"])
        txt_c = t["text"]
    else:  # corporate
        _rect(slide, 0, 0, Inches(0.5), SLIDE_H, t["primary"])
        _rect(slide, Inches(5.5), Inches(4.2), Inches(2.5), Inches(0.04), t["accent"])
        txt_c = t["text"]

    _textbox(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
             title, t["font_title"], Pt(36), txt_c, bold=True,
             align=PP_ALIGN.CENTER)

    if subtitle:
        sub_c = t["text_light"] if tn not in ("finance",) else t["accent"]
        _textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(1),
                 subtitle, t["font_body"], Pt(18), sub_c,
                 align=PP_ALIGN.CENTER)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN BUILD
# ═══════════════════════════════════════════════════════════════════════════════

SLIDE_BUILDERS = {
    "cover": lambda prs, t, tn, d, sn, tot, conf: _slide_cover(prs, t, tn, d),
    "section": lambda prs, t, tn, d, sn, tot, conf: _slide_section(prs, t, tn, d),
    "content": _slide_content,
    "two_col": _slide_two_col,
    "kpi": _slide_kpi,
    "chart": _slide_content,   # fallback to content for now (images ajoutées par le Visualizer)
    "closing": lambda prs, t, tn, d, sn, tot, conf: _slide_closing(prs, t, tn, d),
}


def build(deck: dict, out: Path, theme_override: str = None):
    """Construire un .pptx professionnel à partir d'un dict deck."""
    theme_name = theme_override or deck.get("theme", "corporate")
    if theme_name not in THEMES:
        print(f"[WARN] Thème '{theme_name}' inconnu, fallback → corporate")
        theme_name = "corporate"

    t = THEMES[theme_name]
    confidential = deck.get("confidential", True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides_data = deck.get("slides", [])
    total = len(slides_data)

    for i, sd in enumerate(slides_data):
        slide_type = sd.get("type", "content")

        # Auto-detect cover/closing si non spécifié
        if i == 0 and slide_type == "content":
            slide_type = "cover"
        if i == total - 1 and slide_type == "content" and sd.get("action_title", "").lower() in (
            "merci", "thank you", "questions", "contact", "fin", "closing"):
            slide_type = "closing"

        builder = SLIDE_BUILDERS.get(slide_type, _slide_content)
        builder(prs, t, theme_name, sd, i + 1, total, confidential)

    prs.save(str(out))
    print(f"[OK pptx] {out}  (thème: {theme_name}, {total} slides)")
    return str(out)


def main():
    if len(sys.argv) < 3:
        print("Usage: pptx_builder.py deck.yaml out.pptx [--theme corporate|dark|finance|startup|minimal]")
        sys.exit(1)

    try:
        import yaml
    except ImportError:
        print("[ERR] pip install pyyaml", file=sys.stderr)
        sys.exit(2)

    deck_path = Path(sys.argv[1])
    out_path = Path(sys.argv[2])

    theme = None
    if "--theme" in sys.argv:
        idx = sys.argv.index("--theme")
        if idx + 1 < len(sys.argv):
            theme = sys.argv[idx + 1]

    deck = yaml.safe_load(deck_path.read_text(encoding="utf-8"))
    build(deck, out_path, theme)


if __name__ == "__main__":
    main()
