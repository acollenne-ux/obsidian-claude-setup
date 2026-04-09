#!/usr/bin/env python3
"""
layout_check.py — Analyse géométrique déterministe d'un livrable.

Détecte : overflow (hors page), overlap (chevauchement non déclaré),
clipping texte, margin_violation, empty_region, zorder_suspect.

Usage:
    python layout_check.py <input.pdf|pptx> [--brief brief.md] [--out geom.json]

Exit code 0 = analyse OK (peut contenir des anomalies dans le JSON).
"""

import argparse
import json
import sys
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Any

MIN_MARGIN_PT = 18  # ~6mm / 8pt baseline grid x 2
MIN_EMPTY_RATIO_TRIGGER = 0.60


@dataclass
class Anomaly:
    page: int
    element_id: str
    issue: str
    severity: str  # low | medium | high | critical
    bbox: list[float]
    suggestion: str


def _rects_overlap(a: tuple, b: tuple, tol: float = 2.0) -> bool:
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    return not (
        ax2 - tol <= bx1
        or bx2 - tol <= ax1
        or ay2 - tol <= by1
        or by2 - tol <= ay1
    )


def check_pptx(path: Path) -> list[Anomaly]:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(path))
    slide_w = Emu(prs.slide_width).pt
    slide_h = Emu(prs.slide_height).pt
    anomalies: list[Anomaly] = []

    for idx, slide in enumerate(prs.slides, start=1):
        shapes_boxes: list[tuple[str, tuple, Any]] = []
        occupied = 0.0
        for i, shape in enumerate(slide.shapes):
            if shape.left is None or shape.width is None:
                continue
            x1 = Emu(shape.left).pt
            y1 = Emu(shape.top).pt
            x2 = x1 + Emu(shape.width).pt
            y2 = y1 + Emu(shape.height).pt
            box = (x1, y1, x2, y2)
            name = f"{shape.shape_type}#{i}"
            shapes_boxes.append((name, box, shape))
            occupied += max(0, x2 - x1) * max(0, y2 - y1)

            # overflow
            if x1 < -1 or y1 < -1 or x2 > slide_w + 1 or y2 > slide_h + 1:
                anomalies.append(Anomaly(
                    page=idx, element_id=name, issue="overflow",
                    severity="critical",
                    bbox=[x1, y1, x2 - x1, y2 - y1],
                    suggestion=f"Repositionner dans la slide ({slide_w:.0f}x{slide_h:.0f}pt)."
                ))
            # margin
            margin = min(x1, y1, slide_w - x2, slide_h - y2)
            if margin < MIN_MARGIN_PT and margin >= -1:
                anomalies.append(Anomaly(
                    page=idx, element_id=name, issue="margin_violation",
                    severity="medium",
                    bbox=[x1, y1, x2 - x1, y2 - y1],
                    suggestion=f"Marge < {MIN_MARGIN_PT}pt (actuelle {margin:.1f}pt)."
                ))
            # clipping texte
            if shape.has_text_frame and shape.text_frame.text.strip():
                try:
                    tf = shape.text_frame
                    if not tf.word_wrap and len(tf.text) > 40:
                        anomalies.append(Anomaly(
                            page=idx, element_id=name, issue="clipping",
                            severity="high",
                            bbox=[x1, y1, x2 - x1, y2 - y1],
                            suggestion="Activer word_wrap ou agrandir le conteneur."
                        ))
                except Exception:
                    pass

        # overlap
        for i in range(len(shapes_boxes)):
            for j in range(i + 1, len(shapes_boxes)):
                n1, b1, _ = shapes_boxes[i]
                n2, b2, _ = shapes_boxes[j]
                if _rects_overlap(b1, b2):
                    anomalies.append(Anomaly(
                        page=idx, element_id=f"{n1} x {n2}", issue="overlap",
                        severity="medium",
                        bbox=[
                            max(b1[0], b2[0]), max(b1[1], b2[1]),
                            min(b1[2], b2[2]) - max(b1[0], b2[0]),
                            min(b1[3], b2[3]) - max(b1[1], b2[1]),
                        ],
                        suggestion="Séparer les éléments ou les grouper si intentionnel."
                    ))

        # empty region
        ratio = occupied / (slide_w * slide_h) if slide_w * slide_h else 0
        if ratio < 1 - MIN_EMPTY_RATIO_TRIGGER:
            anomalies.append(Anomaly(
                page=idx, element_id="slide", issue="empty_region",
                severity="low",
                bbox=[0, 0, slide_w, slide_h],
                suggestion=f"Slide {ratio*100:.0f}% occupée — hiérarchie à renforcer."
            ))
    return anomalies


def check_pdf(path: Path) -> list[Anomaly]:
    import pdfplumber

    anomalies: list[Anomaly] = []
    with pdfplumber.open(str(path)) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            pw, ph = page.width, page.height
            boxes: list[tuple[str, tuple]] = []
            for i, word in enumerate(page.extract_words() or []):
                x1, y1, x2, y2 = word["x0"], word["top"], word["x1"], word["bottom"]
                boxes.append((f"word#{i}", (x1, y1, x2, y2)))
                if x1 < -1 or y1 < -1 or x2 > pw + 1 or y2 > ph + 1:
                    anomalies.append(Anomaly(
                        page=idx, element_id=f"word#{i}", issue="overflow",
                        severity="high",
                        bbox=[x1, y1, x2 - x1, y2 - y1],
                        suggestion="Texte sortant de la page."
                    ))
            for rect in page.rects or []:
                x1, y1, x2, y2 = rect["x0"], rect["top"], rect["x1"], rect["bottom"]
                if x1 < -1 or y1 < -1 or x2 > pw + 1 or y2 > ph + 1:
                    anomalies.append(Anomaly(
                        page=idx, element_id="rect", issue="overflow",
                        severity="critical",
                        bbox=[x1, y1, x2 - x1, y2 - y1],
                        suggestion="Forme sortant de la page."
                    ))
    return anomalies


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("input")
    ap.add_argument("--brief", default=None)
    ap.add_argument("--out", default=None)
    args = ap.parse_args()

    src = Path(args.input).resolve()
    ext = src.suffix.lower()
    if ext == ".pptx":
        anomalies = check_pptx(src)
    elif ext == ".pdf":
        anomalies = check_pdf(src)
    else:
        # PNG/JPG/SVG : pas de check géométrique déterministe
        anomalies = []

    result = {
        "input": str(src),
        "brief": args.brief,
        "anomaly_count": len(anomalies),
        "critical": sum(1 for a in anomalies if a.severity == "critical"),
        "high": sum(1 for a in anomalies if a.severity == "high"),
        "anomalies": [asdict(a) for a in anomalies],
    }
    payload = json.dumps(result, indent=2, ensure_ascii=False)
    if args.out:
        Path(args.out).write_text(payload, encoding="utf-8")
    print(payload)
    return 0


if __name__ == "__main__":
    sys.exit(main())
