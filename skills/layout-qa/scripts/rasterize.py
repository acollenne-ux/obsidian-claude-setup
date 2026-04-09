#!/usr/bin/env python3
"""
rasterize.py — Convertit un livrable (PDF/PPTX/PNG/JPG/SVG) en PNG haute résolution.

Usage:
    python rasterize.py <input> --out <dir> [--dpi 200]

Backends:
    - PDF  : pdftoppm (poppler) si dispo, sinon Playwright
    - PPTX : rendu HTML mock via python-pptx → Playwright
    - PNG/JPG : copie
    - SVG  : Playwright
    - LibreOffice INTERDIT (cf. memory 2026-04-07)
"""

import argparse
import json
import shutil
import subprocess
import sys
import uuid
from pathlib import Path


def rasterize_pdf(src: Path, out_dir: Path, dpi: int) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    prefix = out_dir / "page"
    # Try pdftoppm first (fast, deterministic)
    if shutil.which("pdftoppm"):
        subprocess.run(
            ["pdftoppm", "-png", "-r", str(dpi), str(src), str(prefix)],
            check=True,
        )
        return sorted(out_dir.glob("page-*.png"))
    # Fallback: Playwright via pdf.js
    return _rasterize_pdf_playwright(src, out_dir, dpi)


def _rasterize_pdf_playwright(src: Path, out_dir: Path, dpi: int) -> list[Path]:
    from playwright.sync_api import sync_playwright

    html = f"""
    <html><body style="margin:0">
    <embed src="file:///{src.as_posix()}" type="application/pdf"
           width="100%" height="100%"/>
    </body></html>
    """
    tmp_html = out_dir / "_view.html"
    tmp_html.write_text(html, encoding="utf-8")
    outputs: list[Path] = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1600, "height": 2200},
                                device_scale_factor=dpi / 96)
        page.goto(tmp_html.as_uri())
        page.wait_for_timeout(1500)
        out = out_dir / "page-1.png"
        page.screenshot(path=str(out), full_page=True)
        outputs.append(out)
        browser.close()
    return outputs


def rasterize_pptx(src: Path, out_dir: Path, dpi: int) -> list[Path]:
    """Rendu HTML mock de chaque slide via python-pptx → Playwright."""
    from pptx import Presentation
    from pptx.util import Emu
    from playwright.sync_api import sync_playwright

    out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(src))
    slide_w_px = int(Emu(prs.slide_width).pt * dpi / 72)
    slide_h_px = int(Emu(prs.slide_height).pt * dpi / 72)

    outputs: list[Path] = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        for idx, slide in enumerate(prs.slides, start=1):
            html_parts = [
                f'<html><body style="margin:0;width:{slide_w_px}px;'
                f'height:{slide_h_px}px;position:relative;'
                f'font-family:Inter,Arial,sans-serif;background:#fff;">'
            ]
            for shape in slide.shapes:
                try:
                    left = int(Emu(shape.left).pt * dpi / 72) if shape.left else 0
                    top = int(Emu(shape.top).pt * dpi / 72) if shape.top else 0
                    w = int(Emu(shape.width).pt * dpi / 72) if shape.width else 0
                    h = int(Emu(shape.height).pt * dpi / 72) if shape.height else 0
                    text = ""
                    if shape.has_text_frame:
                        text = shape.text_frame.text.replace("\n", "<br>")
                    html_parts.append(
                        f'<div style="position:absolute;left:{left}px;top:{top}px;'
                        f'width:{w}px;height:{h}px;border:1px solid #e5e7eb;'
                        f'padding:6px;box-sizing:border-box;overflow:hidden;'
                        f'font-size:{max(10, h // 20)}px;">{text}</div>'
                    )
                except Exception:
                    continue
            html_parts.append("</body></html>")
            tmp_html = out_dir / f"_slide_{idx}.html"
            tmp_html.write_text("".join(html_parts), encoding="utf-8")

            page = browser.new_page(
                viewport={"width": slide_w_px, "height": slide_h_px}
            )
            page.goto(tmp_html.as_uri())
            page.wait_for_timeout(300)
            out = out_dir / f"page-{idx:03d}.png"
            page.screenshot(path=str(out), full_page=False)
            outputs.append(out)
            page.close()
        browser.close()
    return outputs


def rasterize_image(src: Path, out_dir: Path) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    dst = out_dir / f"page-001{src.suffix.lower()}"
    shutil.copy2(src, dst)
    return [dst]


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("input")
    ap.add_argument("--out", required=True)
    ap.add_argument("--dpi", type=int, default=200)
    args = ap.parse_args()

    src = Path(args.input).resolve()
    out_dir = Path(args.out).resolve() / uuid.uuid4().hex[:8]
    ext = src.suffix.lower()

    if ext == ".pdf":
        pages = rasterize_pdf(src, out_dir, args.dpi)
    elif ext == ".pptx":
        pages = rasterize_pptx(src, out_dir, args.dpi)
    elif ext in (".png", ".jpg", ".jpeg", ".svg"):
        pages = rasterize_image(src, out_dir)
    else:
        print(f"ERROR: unsupported format {ext}", file=sys.stderr)
        return 2

    print(json.dumps({"out_dir": str(out_dir),
                      "pages": [str(p) for p in pages]}, indent=2))
    return 0


if __name__ == "__main__":
    sys.exit(main())
